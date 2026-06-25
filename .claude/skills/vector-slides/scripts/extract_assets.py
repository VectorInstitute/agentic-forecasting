#!/usr/bin/env python3
"""
Extract the Vector brand asset library from the source decks.

This is a **lab/dev tool**, not part of the consumer build path. It pulls the
full-resolution visual elements the composable compiler needs — the gradient
hero background, the "A" arrow overlay, the flat brand-colored icon set, and the
embedded Open Sans / Montserrat fonts — out of the two source .pptx files and
writes them into ``assets/brand/`` next to this skill, plus a ``brand.json``
manifest.

Once extracted and committed, the published skill ships only these small assets
(~a few MB) instead of the two ~57 MB master templates.

Sources (defaults, overridable via CLI):
  - Master Presentation Template.pptx  → embedded fonts, sector photos
  - full-deck.pptx                     → gradient bg, arrow, icon set

Usage:
  uv run python scripts/extract_assets.py \
      --master "/path/Master Presentation Template.pptx" \
      --deck   "/path/full-deck.pptx" \
      --out    assets/brand
"""
from __future__ import annotations

import argparse
import json
import shutil
import zipfile
from pathlib import Path

# full-deck.pptx media part  ->  semantic icon name (+ default baked color).
# These part names are stable in the source deck; re-verify if the deck changes
# (run with --report to dump the current picture→media map).
ICON_MAP: dict[str, dict[str, str]] = {
    "image30.png": {"name": "flask", "color": "purple"},
    "image19.png": {"name": "shield", "color": "pink"},
    "image1.png": {"name": "book", "color": "blue"},
    "image8.png": {"name": "bug", "color": "red"},
    "image18.png": {"name": "brain", "color": "pink"},
    "image22.png": {"name": "warning", "color": "amber"},
    "image16.png": {"name": "check", "color": "green"},
    "image4.png": {"name": "x", "color": "red"},
    "image6.png": {"name": "arrow", "color": "pink"},
    "image9.png": {"name": "chart", "color": "pink"},
    "image25.png": {"name": "code", "color": "blue"},
    "image34.png": {"name": "search", "color": "blue"},
    "image37.png": {"name": "gear", "color": "gray"},
    "image43.png": {"name": "robot", "color": "blue"},
}

# Hero assets from the MASTER template (the real brand look): soft mesh gradients,
# the rising gradient arrow, and the Vector logo lockups.
MASTER_BACKGROUNDS = {
    "image10.png": "mesh-magenta",  # pink/magenta ↔ blue (title default)
    "image14.png": "mesh-lime",     # lime/yellow ↔ blue
    "image15.png": "mesh-cyan",     # cyan/teal ↔ violet
    "image8.png": "mesh-amber",     # amber/orange ↔ purple
}
MASTER_ARROW = "image29.png"       # pink→purple→blue rising arrow (signature motif)
MASTER_LOGO_WHITE = "image1.png"   # horizontal white lockup (for dark/gradient heroes)
MASTER_LOGO_MARK = "image6.png"    # stacked black "V"+pink arrow + wordmark (for white footers)

# Brand palette (official sampled values — NOT the generic Office theme).
PALETTE = {
    "pink": "FF008C",
    "blue": "313CFF",
    "purple": "8A25C9",
    "cyan": "48C0D9",
    "amber": "FF9E00",
    "lime": "CFF933",
    "green": "1DB47F",   # success accent used in compare/check
    "red": "E8553A",     # warning/wrong accent used in compare/x
    "black": "000000",
    "white": "FFFFFF",
    "ink": "222222",     # primary body text on light
    "gray": "666666",    # secondary body text
    "card": "F5F5F5",    # card fill
}

FONTS = {
    # role -> (family name, source .fntdata part basename in the master)
    "heading": "Open Sans",
    "body": "Open Sans",
    "display": "Open Sans",
}


def _extract_part(zf: zipfile.ZipFile, member: str, dest: Path) -> bool:
    try:
        with zf.open(member) as src, open(dest, "wb") as out:
            shutil.copyfileobj(src, out)
        return True
    except KeyError:
        return False


def extract(master: Path, deck: Path, out: Path) -> dict:
    icons_dir = out / "icons"
    bg_dir = out / "backgrounds"
    fonts_dir = out / "fonts"
    for d in (icons_dir, bg_dir, fonts_dir):
        d.mkdir(parents=True, exist_ok=True)

    manifest: dict = {
        "canvas": {"width_in": 10.0, "height_in": 5.625, "aspect": "16:9"},
        "palette": PALETTE,
        "fonts": FONTS,
        "footer": {
            "bar_top_in": 5.20, "bar_height_in": 0.42,
            "rule_height_in": 0.03, "rule_color": "pink",
            "label": "VECTOR INSTITUTE", "label_color": "white",
        },
        "backgrounds": {}, "arrow": None, "logos": {}, "icons": {}, "font_files": [],
    }

    # Icons come from full-deck (clean flat brand-colored set).
    with zipfile.ZipFile(deck) as zf:
        for media, meta in ICON_MAP.items():
            dest = icons_dir / f"{meta['name']}.png"
            if _extract_part(zf, f"ppt/media/{media}", dest):
                manifest["icons"][meta["name"]] = {
                    "path": f"icons/{meta['name']}.png", "color": meta["color"],
                }

    # Hero backgrounds, rising arrow, logos, and fonts come from the MASTER.
    with zipfile.ZipFile(master) as zf:
        for media, name in MASTER_BACKGROUNDS.items():
            if _extract_part(zf, f"ppt/media/{media}", bg_dir / f"{name}.png"):
                manifest["backgrounds"][name] = f"backgrounds/{name}.png"
        if _extract_part(zf, f"ppt/media/{MASTER_ARROW}", out / "arrow-rising.png"):
            manifest["arrow"] = "arrow-rising.png"
        if _extract_part(zf, f"ppt/media/{MASTER_LOGO_WHITE}", out / "logo-white.png"):
            manifest["logos"]["white"] = "logo-white.png"
        if _extract_part(zf, f"ppt/media/{MASTER_LOGO_MARK}", out / "logo-mark.png"):
            manifest["logos"]["mark"] = "logo-mark.png"
        for member in zf.namelist():
            if member.startswith("ppt/fonts/") and member.endswith(".fntdata"):
                name = Path(member).name
                if _extract_part(zf, member, fonts_dir / name):
                    manifest["font_files"].append(f"fonts/{name}")

    (out / "brand.json").write_text(json.dumps(manifest, indent=2) + "\n")
    return manifest


def report(deck: Path) -> None:
    """Dump every picture's media part + size to help re-map icons if the deck changes."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from pptx.util import Emu

    prs = Presentation(str(deck))
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not str(sh.shape_type).startswith("PICTURE"):
                continue
            blip = sh._element.find(".//" + qn("a:blip"))
            rid = blip.get(qn("r:embed")) if blip is not None else None
            part = sh.part.related_part(rid).partname if rid else "?"
            w = Emu(sh.width).inches
            print(f"  s{i:02d} {part} {w:.2f}in")


def main() -> int:
    root = Path(__file__).resolve().parent.parent  # skill root
    lab_assets = root.parent.parent.parent / "assets"  # lab repo assets/ (best-effort default)
    ap = argparse.ArgumentParser(description="Extract Vector brand asset library")
    ap.add_argument("--master", type=Path, default=lab_assets / "Master Presentation Template.pptx")
    ap.add_argument("--deck", type=Path, default=lab_assets / "full-deck.pptx")
    ap.add_argument("--out", type=Path, default=root / "assets" / "brand")
    ap.add_argument("--report", action="store_true", help="dump picture→media map and exit")
    args = ap.parse_args()

    if args.report:
        report(args.deck)
        return 0

    for p in (args.master, args.deck):
        if not p.exists():
            print(f"Source not found: {p}")
            return 1

    import shutil as _sh
    if args.out.exists():
        _sh.rmtree(args.out)  # clean stale assets from prior extractions
    manifest = extract(args.master, args.deck, args.out)
    print(f"Extracted brand assets to {args.out}")
    print(f"  backgrounds ({len(manifest['backgrounds'])}): {', '.join(manifest['backgrounds'])}")
    print(f"  arrow: {manifest['arrow']}  logos: {list(manifest['logos'])}")
    print(f"  icons ({len(manifest['icons'])}): {', '.join(sorted(manifest['icons']))}")
    print(f"  font files: {len(manifest['font_files'])}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
