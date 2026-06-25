"""
Composable slide primitives — the building blocks of the Vector deck compiler.

Every primitive emits **native python-pptx shapes** (auto-shapes, text boxes,
pictures, tables) which serialize to valid OOXML. We never hand-build or inject
raw ``<p:sp>`` XML strings (that is what caused the historical PowerPoint
"repair" prompts). Geometry and color defaults come from ``brand.py``.

Coordinates are in inches (floats). Colors are palette names or 6-digit hex.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import brand
from brand import FONT_BODY, FONT_HEADING, color

# ---------------------------------------------------------------------------
# Slide + background
# ---------------------------------------------------------------------------
def blank_slide(prs):
    """Add a slide on the blank DEFAULT layout."""
    layout = prs.slide_layouts[0]  # base has exactly one blank layout
    return prs.slides.add_slide(layout)


def set_bg(slide, fill="white") -> None:
    """Set a solid slide background color (valid <p:bg> on the slide)."""
    cSld = slide._element.find(qn("p:cSld"))
    # remove any existing bg
    existing = cSld.find(qn("p:bg"))
    if existing is not None:
        cSld.remove(existing)
    bg = cSld.makeelement(qn("p:bg"), {})
    bgPr = bg.makeelement(qn("p:bgPr"), {})
    solidFill = bgPr.makeelement(qn("a:solidFill"), {})
    srgb = solidFill.makeelement(qn("a:srgbClr"), {"val": str(color(fill))})
    solidFill.append(srgb)
    bgPr.append(solidFill)
    bgPr.append(bgPr.makeelement(qn("a:effectLst"), {}))
    bg.append(bgPr)
    cSld.insert(0, bg)  # bg must be first child of cSld


def full_bleed_image(slide, path: Path):
    """Add an image covering the whole canvas (added first → sits at back)."""
    return slide.shapes.add_picture(
        str(path), Inches(0), Inches(0), brand.CANVAS_W, brand.CANVAS_H
    )


def image(slide, path: Path, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kw)


def image_fit(slide, path: Path, x, y, max_w, max_h, *, halign="center",
              valign="middle"):
    """Place a picture scaled to fit inside ``(max_w, max_h)``, preserving aspect.

    Returns the picture. Used by the figure/figure_full layouts so a plot of any
    aspect ratio sits cleanly in its content slot without distortion.
    """
    from PIL import Image as _Img

    iw, ih = _Img.open(str(path)).size
    ar = iw / ih if ih else 1.0
    w = max_w
    h = w / ar
    if h > max_h:
        h = max_h
        w = h * ar
    if halign == "center":
        x = x + (max_w - w) / 2
    elif halign == "right":
        x = x + (max_w - w)
    if valign == "middle":
        y = y + (max_h - h) / 2
    elif valign == "bottom":
        y = y + (max_h - h)
    return image(slide, path, x, y, w=w, h=h)


def rounded_picture(slide, path: Path, x, y, w, h, radius_in=0.22):
    """Add a picture with rounded corners (swap its rect geometry for roundRect)."""
    pic = image(slide, path, x, y, w=w, h=h)
    spPr = pic._element.spPr
    old = spPr.find(qn("a:prstGeom"))
    if old is not None:
        spPr.remove(old)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "roundRect"})
    av = geom.makeelement(qn("a:avLst"), {})
    adj = max(0.0, min(0.5, radius_in / min(w, h)))
    gd = av.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {int(adj * 100000)}"})
    av.append(gd)
    geom.append(av)
    # geometry must precede fill/line in spPr; insert after xfrm if present
    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is not None:
        xfrm.addnext(geom)
    else:
        spPr.insert(0, geom)
    return pic


def hero_panel(slide, variant=None, *, framed=True):
    """White slide + a rounded mesh-gradient panel (the master title/section look)."""
    set_bg(slide, "white")
    bg = brand.mesh_bg(variant)
    if framed:
        m = 0.13
        return rounded_picture(slide, bg, m, m, brand.CANVAS_W_IN - 2 * m,
                               brand.CANVAS_H_IN - 2 * m, radius_in=0.26)
    return full_bleed_image(slide, bg)


def hero_logo(slide):
    """White Vector lockup, bottom-left, on a hero panel."""
    if not brand.LOGO_WHITE.exists():
        return
    from PIL import Image as _Img
    iw, ih = _Img.open(brand.LOGO_WHITE).size
    w = 2.0
    h = w * ih / iw
    image(slide, brand.LOGO_WHITE, 0.62, brand.CANVAS_H_IN - 0.45 - h, w=w, h=h)


def rising_arrow(slide, *, w=3.0, right=0.4, bottom=0.0):
    """The pink→purple rising gradient arrow, anchored bottom-right (signature motif)."""
    if not brand.ARROW_RISING.exists():
        return
    from PIL import Image as _Img
    iw, ih = _Img.open(brand.ARROW_RISING).size
    h = w * ih / iw
    x = brand.CANVAS_W_IN - right - w
    y = brand.CANVAS_H_IN - bottom - h
    image(slide, brand.ARROW_RISING, x, y, w=w, h=h)


def _gradient_fill_elem(holder, c1, c2, angle_deg=45):
    """Build + attach an <a:gradFill> (2-stop linear) onto an spPr or ln element."""
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        e = holder.find(qn(tag))
        if e is not None:
            holder.remove(e)
    grad = holder.makeelement(qn("a:gradFill"), {})
    lst = grad.makeelement(qn("a:gsLst"), {})
    for pos, c in ((0, c1), (100000, c2)):
        gs = lst.makeelement(qn("a:gs"), {"pos": str(pos)})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": str(color(c))})
        gs.append(clr)
        lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(int(angle_deg * 60000)), "scaled": "1"})
    grad.append(lin)
    holder.append(grad)
    return grad


def gradient_box(slide, x, y, w, h, pair, *, border_pt=4.0, fill="white", radius_in=0.12):
    """Rounded rect with a thick gradient border + light fill (master 'colour box')."""
    c1, c2 = pair
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shape.adjustments[0] = max(0.0, min(0.5, radius_in / min(w, h)))
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.shadow.inherit = False
    # gradient line
    ln = shape._element.spPr.get_or_add_ln()
    ln.set("w", str(int(border_pt * 12700)))
    _gradient_fill_elem(ln, c1, c2, angle_deg=30)
    return shape


# ---------------------------------------------------------------------------
# Rectangles / bars / cards
# ---------------------------------------------------------------------------
def _no_line(shape) -> None:
    shape.line.fill.background()


def _soft_shadow(shape) -> None:
    """Attach a subtle outer shadow (native a:effectLst)."""
    spPr = shape._element.spPr
    # remove inherited effect refs / lists first
    for tag in ("a:effectLst", "a:effectDag"):
        e = spPr.find(qn(tag))
        if e is not None:
            spPr.remove(e)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shd = eff.makeelement(
        qn("a:outerShdw"),
        {"blurRad": "50800", "dist": "25400", "dir": "5400000", "rotWithShape": "0"},
    )
    clr = shd.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "14000"}))
    shd.append(clr)
    eff.append(shd)
    spPr.append(eff)


def rect(slide, x, y, w, h, *, fill=None, line=None, line_w=1.0,
         rounded=False, radius_in=0.08, shadow=False):
    """A rectangle (optionally rounded). fill/line are palette names or hex, or None."""
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            adj = max(0.0, min(0.5, radius_in / min(w, h)))
            shape.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color(fill)
    if line is None:
        _no_line(shape)
    else:
        shape.line.color.rgb = color(line)
        shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    if shadow:
        _soft_shadow(shape)
    return shape


def card(slide, x, y, w, h, *, fill=brand.CARD_FILL, accent=None, accent_side="top",
         shadow=True, rounded=True):
    """A content card with an optional colored accent bar on one edge."""
    bg = rect(slide, x, y, w, h, fill=fill, rounded=rounded, shadow=shadow)
    if accent:
        if accent_side == "top":
            rect(slide, x, y, w, brand.ACCENT_H, fill=accent)
        elif accent_side == "left":
            rect(slide, x, y, brand.ACCENT_W, h, fill=accent)
    return bg


def accent_bar(slide, x, y, w, h, fill="pink"):
    return rect(slide, x, y, w, h, fill=fill)


# ---------------------------------------------------------------------------
# Text
# ---------------------------------------------------------------------------
_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
_ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def _suppress_bullet(para) -> None:
    """Force no bullet on a paragraph (some inherited styles add one)."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        for e in pPr.findall(qn(tag)):
            pPr.remove(e)
    # insert <a:buNone/> before defRPr/extLst if present, else append
    bunone = pPr.makeelement(qn("a:buNone"), {})
    ref = None
    for tag in ("a:defRPr", "a:extLst"):
        ref = pPr.find(qn(tag))
        if ref is not None:
            break
    if ref is not None:
        ref.addprevious(bunone)
    else:
        pPr.append(bunone)


def _style_run(run, *, font, size, bold, italic, rgb) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb
    # pin latin typeface so it survives in PowerPoint
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {})
        rPr.append(latin)
    latin.set("typeface", font)


def add_text(slide, x, y, w, h, content, *, size=brand.TYPE["body"], color="ink",
             bold=False, italic=False, font=FONT_BODY, align="left", anchor="top",
             line_spacing=1.0, space_after=4.0, wrap=True):
    """
    Add a text box.

    ``content`` may be:
      - str                          → one paragraph
      - list[str]                    → one paragraph per item
      - list[dict] paragraphs, each: {text|runs, size?, color?, bold?, italic?, align?, space_after?}
        where runs = [{text, bold?, italic?, color?, size?}, ...]
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = _ANCHOR[anchor]
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)

    paragraphs = _normalize_content(content)
    base = dict(size=size, color=color, bold=bold, italic=italic, font=font,
                align=align, line_spacing=line_spacing, space_after=space_after)
    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _suppress_bullet(para)
        pstyle = {**base, **{k: v for k, v in p.items() if k != "text" and k != "runs"}}
        para.alignment = _ALIGN[pstyle["align"]]
        if pstyle.get("line_spacing"):
            para.line_spacing = pstyle["line_spacing"]
        para.space_before = Pt(0)
        para.space_after = Pt(pstyle.get("space_after", 0))
        runs = p.get("runs")
        if runs is None:
            runs = [{"text": p.get("text", "")}]
        for r in runs:
            run = para.add_run()
            run.text = r.get("text", "")
            _style_run(
                run,
                font=r.get("font", pstyle["font"]),
                size=r.get("size", pstyle["size"]),
                bold=r.get("bold", pstyle["bold"]),
                italic=r.get("italic", pstyle["italic"]),
                rgb=color_fn(r.get("color", pstyle["color"])),
            )
    return box


def color_fn(c) -> RGBColor:
    return c if isinstance(c, RGBColor) else color(c)


def _normalize_content(content) -> list[dict]:
    if content is None:
        return [{"text": ""}]
    if isinstance(content, str):
        return [{"text": content}]
    out = []
    for item in content:
        if isinstance(item, str):
            out.append({"text": item})
        elif isinstance(item, dict):
            out.append(item)
        else:
            out.append({"text": str(item)})
    return out


# ---------------------------------------------------------------------------
# Composite primitives
# ---------------------------------------------------------------------------
def icon(slide, name: str, x, y, size=0.34):
    return image(slide, brand.icon_path(name), x, y, w=size, h=size)


def number_circle(slide, x, y, d, n, *, fill="pink", text_color="white",
                  size=brand.TYPE["card_title"]):
    """A filled circle with a centered number."""
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    oval.fill.solid()
    oval.fill.fore_color.rgb = color(fill)
    _no_line(oval)
    oval.shadow.inherit = False
    tf = oval.text_frame
    tf.word_wrap = False
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    _suppress_bullet(para)
    run = para.add_run()
    run.text = str(n)
    _style_run(run, font=FONT_HEADING, size=size, bold=True, italic=False,
               rgb=color(text_color))
    return oval


def pill(slide, x, y, w, h, text, *, fill="white", text_color="ink",
         size=brand.TYPE["body"], bold=False, align="center"):
    shape = rect(slide, x, y, w, h, fill=fill, rounded=True, radius_in=0.05)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    para = tf.paragraphs[0]
    para.alignment = _ALIGN[align]
    _suppress_bullet(para)
    run = para.add_run()
    run.text = text
    _style_run(run, font=FONT_BODY, size=size, bold=bold, italic=False,
               rgb=color(text_color))
    return shape


def callout_bar(slide, x, y, w, text, *, h=0.45, fill="black", text_color="pink",
                size=brand.TYPE["callout"], align="center", bold=True, italic=False):
    bar = rect(slide, x, y, w, h, fill=fill)
    add_text(slide, x + 0.25, y, w - 0.5, h, text, size=size, color=text_color,
             bold=bold, italic=italic, align=align, anchor="middle", space_after=0)
    return bar


def footer(slide):
    """Master footer: a thin pink rule + the stacked Vector logo lockup, bottom-right."""
    f = brand.FOOTER
    rect(slide, f["rule_left"], f["rule_y"], f["rule_w"], f["rule_h"], fill=f["rule_color"])
    if brand.LOGO_MARK.exists():
        from PIL import Image as _Img
        iw, ih = _Img.open(brand.LOGO_MARK).size
        w = f["logo_w"]
        h = w * ih / iw
        x = brand.CANVAS_W_IN - f["logo_right_pad"] - w
        y = brand.CANVAS_H_IN - f["logo_bottom_pad"] - h
        image(slide, brand.LOGO_MARK, x, y, w=w, h=h)


# ---------------------------------------------------------------------------
# Title / subtitle helpers (content slides)
# ---------------------------------------------------------------------------
def content_title(slide, text, *, top=brand.TITLE_TOP, size=brand.TYPE["title"]):
    add_text(slide, brand.MARGIN_X, top, brand.CONTENT_W, brand.TITLE_H, text,
             size=size, color="ink", bold=True, font=FONT_HEADING, space_after=0)


def content_subtitle(slide, text, *, top=brand.SUBTITLE_TOP):
    if not text:
        return
    add_text(slide, brand.MARGIN_X, top, brand.CONTENT_W, 0.35, text,
             size=brand.TYPE["subtitle"], color="muted", italic=True, space_after=0)


# ---------------------------------------------------------------------------
# Table
# ---------------------------------------------------------------------------
def data_table(slide, x, y, w, h, headers, rows, *, highlights=None):
    """
    Native table. headers: list[str]; rows: list[list[str]].
    highlights: optional dict {(row_idx, col_idx): palette_name} for cell fills,
    where row_idx 0 = first body row.
    """
    nrows = len(rows) + 1
    ncols = len(headers)
    gfx = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = gfx.table
    table.first_row = False
    table.horz_banding = False

    def fill_cell(cell, text, *, bold, txt_color, fill_color, size, align=PP_ALIGN.CENTER):
        cell.fill.solid()
        cell.fill.fore_color.rgb = color(fill_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        _suppress_bullet(para)
        run = para.add_run()
        run.text = str(text)
        _style_run(run, font=FONT_BODY, size=size, bold=bold, italic=False,
                   rgb=color(txt_color))

    for c, head in enumerate(headers):
        fill_cell(table.cell(0, c), head, bold=True, txt_color="white",
                  fill_color="black", size=brand.TYPE["table"],
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    highlights = highlights or {}
    for r, row in enumerate(rows):
        zebra = "white" if r % 2 else "card"
        for c, val in enumerate(row):
            fill_color = highlights.get((r, c), zebra)
            is_label = c == 0
            fill_cell(
                table.cell(r + 1, c), val,
                bold=is_label, txt_color="ink", fill_color=fill_color,
                size=brand.TYPE["table"],
                align=PP_ALIGN.LEFT if is_label else PP_ALIGN.CENTER,
            )
    return table
