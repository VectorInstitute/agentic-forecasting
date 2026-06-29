"""
Geometry-aware text-overflow heuristic for composed decks.

Walks every text-bearing shape (text boxes, auto-shapes, table cells) and
estimates whether the text fits its box, using the run font sizes the compiler
actually set. Composed boxes are sized tightly, so overflow = a real visual
collision (text boxes don't clip — they spill onto whatever is below).

Calibrated to flag only clear overflows (≥1 extra line beyond capacity on a box
that needs ≥2 lines) to avoid false positives on headings and one-line labels.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt

LINE_HEIGHT = 1.22     # line box as a fraction of font size
MARGIN_W = 0.10        # inches lost to internal L+R margins (approx)
MARGIN_H = 0.04


def _advance(font_pt: float) -> float:
    """Avg glyph width as a fraction of em. Large display text is bold and its
    metrics run wider, so big headings get a higher advance — this is what lets
    the check catch a too-long title that wraps into the content below it."""
    return 0.52 if font_pt >= 24 else 0.50


def _max_font_pt(tf, default=12.0) -> float:
    sizes = [r.font.size.pt for p in tf.paragraphs for r in p.runs if r.font.size]
    return max(sizes) if sizes else default


def _lines_needed(tf, cpl: int) -> int:
    total = 0
    for p in tf.paragraphs:
        text = "".join(r.text for r in p.runs)
        total += max(1 if text.strip() else 0, -(-len(text) // cpl))
    return max(total, 1)


# A single line of text needs at least ~0.9·em of box height (cap height + the
# descender) or its glyphs spill past the box's bottom edge. Below this the box is
# simply too short for even one line — a border/edge overlap the line-count check
# misses because it floors capacity to "1 line fits".
GLYPH_FRAC = 0.9


def _flag(w_in, h_in, tf, *, font_pt) -> tuple[int, int] | None:
    usable_w = max(0.2, w_in - MARGIN_W)
    usable_h = max(0.1, h_in - MARGIN_H)
    cpl = max(1, int(usable_w * 72 / (font_pt * _advance(font_pt))))
    fit = max(1, int(usable_h * 72 / (font_pt * LINE_HEIGHT)))
    need = _lines_needed(tf, cpl)
    # Box too short to contain even one line of its own text → glyphs spill past
    # the bottom edge (flagged regardless of line count).
    if need >= 1 and usable_h < (font_pt / 72) * GLYPH_FRAC:
        return need, 0
    if need > fit and (need - fit) >= 1 and need >= 2:
        return need, fit
    return None


def estimate_overflows(pptx_path: Path) -> list[dict]:
    prs = Presentation(str(pptx_path))
    out: list[dict] = []
    for sidx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_table:
                _check_table(shape, sidx, out)
                continue
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text
            if not text.strip():
                continue
            w_in = Emu(shape.width).inches
            h_in = Emu(shape.height).inches
            font_pt = _max_font_pt(tf)
            res = _flag(w_in, h_in, tf, font_pt=font_pt)
            if res:
                need, fit = res
                cpl = max(1, int((w_in - MARGIN_W) * 72 / (font_pt * _advance(font_pt))))
                out.append({
                    "slide": sidx, "kind": "text",
                    "lines_needed": need, "lines_fit": fit, "cpl": cpl,
                    "chars": len(text), "font_pt": round(font_pt, 1),
                    "preview": text.replace("\n", " / ")[:60],
                })
    return out


def _check_table(shape, sidx, out) -> None:
    table = shape.table
    # approximate per-cell width from total table width / col count
    ncols = len(table.columns)
    w_total = Emu(shape.width).inches
    for r, row in enumerate(table.rows):
        h_in = Emu(row.height).inches
        for c in range(ncols):
            cell = table.cell(r, c)
            tf = cell.text_frame
            if not tf.text.strip():
                continue
            font_pt = _max_font_pt(tf)
            res = _flag(w_total / ncols, h_in, tf, font_pt=font_pt)
            if res:
                need, fit = res
                cpl = max(1, int((w_total / ncols - MARGIN_W) * 72 / (font_pt * _advance(font_pt))))
                out.append({
                    "slide": sidx, "kind": "table-cell",
                    "lines_needed": need, "lines_fit": fit, "cpl": cpl,
                    "chars": len(tf.text), "font_pt": round(font_pt, 1),
                    "preview": f"r{r}c{c}: " + tf.text.replace("\n", " ")[:50],
                })
